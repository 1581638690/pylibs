# 正式使用
import sys

sys.path.append("/opt/openfbi/pylibs/")
from id_validator import validator
from pptx import Presentation

# import paddle
# import jieba.posseg as pseg
try:
    import regex
except:
    import re
import bz2
import gzip
import lzma
import shutil
import tarfile
import docx
import filetype
import zipfile
import openpyxl

try:
    import pdfplumber
except:
    pass
from unrar import rarfile
import xlrd
import os



def guess_file(filename, t_path, yuanfile_path,file_set,magic):
    """
    :param filename:  文件路径
    :param name: 文件名
    :param t_path: 文件解压到的路径
    :param yuanfile_path: 文件全局变量存放
    :return:
    """
    zip_list=["zip", "word", "ppt", "excel"]
    # 用来存放第一次进来的类型
    # global name_type
    # error_count = 0
    # max_attempts = 5
    # while error_count <= max_attempts:
    #     try:
    #         name_type = filetype.guess(filename)
    #         # 如果成功获取文件类型判断 则跳出
    #         break
    #     except FileNotFoundError:
    #         name_type = None
    #     except Exception as e:
    #         # 处理其他异常情况
    #         error_count += 1
    #         name_type = None
    #         raise Exception("发生异常", e)
    # if name_type is not None:
    #     file_type = name_type.mime.split("/")[1]
    # else:
    #     file_type = ""
    name = filename.split("/")[-1]
    print(name)
    # print(filename,name)
    #if file_type == "zip":
    if "zip" in magic.lower() and all(ext in zip_list for ext in file_set):
        try:
            zf = zipfile.ZipFile(filename)
            # print(zf.namelist())
            an = check_zip(filename)
            zip_info = {}
            if an:
                yuanfile_path[name] = "加密zip"
            else:
                namelist = zf.namelist()  #
                zf.extractall(t_path)
                if len(namelist) > 1:
                    # 判断是否为excel和word文档
                    word_type = False
                    xl_type = False
                    ppt_type = False
                    if namelist[0] == "[Content_Types].xml":
                        for n in namelist:
                            if "word" in n:
                                word_type = True
                                break
                            elif "xl" in n:
                                xl_type = True
                                break
                            elif "ppt" in n:
                                ppt_type = True
                                break
                        if word_type:
                            content = ""
                            try:
                                content = word_con(filename, content)
                            except:
                                content = ""
                            yuanfile_path[name] = content
                        if xl_type:
                            content = ""
                            try:
                                content = excel_con(zf, filename, content)
                            except:
                                content = ""
                            yuanfile_path[name] = content
                        # 2023/1/13 add by rzc  python-pptx 文件
                        if ppt_type:
                            content = ""
                            try:
                                content = ppt_con(filename, content)
                            except:
                                content = ""
                            yuanfile_path[name] = content
                    else:
                        # 取出地址最后一段
                        un = filename.split("/")[-1]  # TEST.zip
                        # name_file作为文件路径,文件为zip文件且下面包含多个文件
                        # 如果存在解压文件第一个等于文件路径得最后一段，那么再对每一个文件进行解压
                        yuanfile_path[un] = []
                        for filename in namelist[1:]:
                            # 还是调用file_path函数
                            file_path = t_path + filename
                            # print(file_path)
                            try:
                                info = guess_file(file_path, t_path, zip_info)
                                yuanfile_path[un].append(info)
                                zip_info = {}
                            except:
                                pass

                else:
                    for name in namelist:
                        filepath = t_path + name

                        with open(filepath, 'rb') as fp:
                            try:
                                content = fp.read().decode("utf-8")
                                yuanfile_path[name] = content
                            except:
                                content = fp.read().decode("gbk")
                                yuanfile_path[name] = content
        except:
            pass
    elif "rar" in magic.lower() and "rar" in file_set:
        # 判断传入yuanfile_info是否存在键值
        # 判断rar文件
        # 将rar_info作为一个新的变量传入函数
        try:
            rar_info = {}
            rf = rarfile.RarFile(filename)
            namelist = rf.namelist()
            rf.extractall(t_path)
            if len(namelist) > 1:
                # 取出地址最后一段
                un = filename.split("/")[-1]  # TEST.zip
                yuanfile_path[un] = []
                for n in namelist[:-1]:
                    file_path = t_path + n
                    # file_path=file_path.replace("/","\\")
                    try:
                        info = guess_file(file_path, t_path, rar_info)  # 递归调用
                        yuanfile_path[un].append(info)
                        rar_info = {}
                    except:
                        pass
            else:
                for name in namelist:
                    filepath = t_path + name
                    # filepath1 = filepath.replace("/", "\\")
                    with open(filepath, 'rb') as fp:
                        try:
                            content = fp.read().decode("utf-8")
                            yuanfile_path[name] = content
                        except:
                            content = fp.read().decode("gbk")
                            yuanfile_path[name] = content
        except:
            pass
    elif "tar" in magic.lower() and "tar" in file_set:
        # 将tar_info作为一个新的变量传入函数
        try:
            tar_info = {}
            # 判断tar文件
            tf = tarfile.open(filename)
            namelist = tf.getnames()
            tf.extractall(t_path)
            # 判断namelist是否大于1
            if len(namelist) > 1:
                # 取出地址最后一段
                un = filename.split("/")[-1]  # TEST.zip
                # name_file作为文件路径,文件为zip文件且下面包含多个文件
                yuanfile_path[un] = []
                for name in namelist[1:]:
                    file_path = t_path + name
                    # file_path1=file_path.replace("/","\\")
                    try:
                        info = guess_file(file_path, t_path, tar_info)  # 递归调用
                        yuanfile_path[un].append(info)
                        tar_info = {}
                    except:
                        pass
            else:
                for name in namelist:
                    filepath = t_path + name
                    # filepath1 = filepath.replace("/", "\\")
                    with open(filepath, 'rb') as fp:
                        try:
                            content = fp.read().decode("utf-8")
                            yuanfile_path[name] = content
                        except:
                            content = fp.read().decode("gbk")
                            yuanfile_path[name] = content
        except:
            pass
    elif "gz" in magic.lower() and "gz" in file_set:
        an = gzip.open(filename, mode='r')
        try:
            content = an.read().decode("utf-8")
            yuanfile_path[name] = content
        except:
            try:
                content = an.read().decode("gbk")
                yuanfile_path[name] = content
            except:
                pass
    elif "x-bzip2" in magic.lower() and "x-bzip2" in file_set:
        an = bz2.BZ2File(filename)
        try:
            content = an.read().decode("utf-8")
            yuanfile_path[name] = content
        except:
            try:
                content = an.read().decode("gbk")
                yuanfile_path[name] = content
            except:
                pass
        # print(yuanfile_path)
    elif "xz" in magic.lower() and "xz" in file_set:
        readfile = lzma.open(filename, 'rb')
        try:
            content = readfile.read().decode("utf-8")
            yuanfile_path[name] = content
        except:
            try:
                content = readfile.read().decode("gbk")
                yuanfile_path[name] = content
            except:
                pass

    elif "pdf" in magic.lower() and "pdf" in file_set:
        try:
            content = ""
            with pdfplumber.open(filename) as fp:
                page_list = fp.pages
                for page in page_list:
                    content += page.extract_text()
            yuanfile_path[name] = content
        except:
            pass
    else:

        with open(filename, 'rb') as fp:
            try:
                content = fp.read().decode("utf-8")
                yuanfile_path[name] = content
            except:
                try:
                    content = fp.read().decode("gbk")
                    yuanfile_path[name] = content
                except:
                    pass

    return yuanfile_path


def check_zip(filename):
    zf = zipfile.ZipFile(filename)
    for zinfo in zf.infolist():
        is_encrypted = zinfo.flag_bits & 0x1
        if is_encrypted:
            return True
        else:
            return False


# 删除某一目录下的所有文件或文件夹
def del_file(t_path):
    del_list = os.listdir(t_path)
    for f in del_list:
        file_path = os.path.join(t_path, f)
        # 2023-8-9 添加链接符号判断 避免出现"Cannot call rmtree on a symbolic link"错误
        if os.path.islink(file_path):
            os.unlink(file_path)
        elif os.path.isfile(file_path):
            os.remove(file_path)
        elif os.path.isdir(file_path):
            shutil.rmtree(file_path)


# 获取文件大小
def getDocSize(path):
    try:
        size = os.path.getsize(path)
        return size
    except Exception as err:
        print(err)


# 对获取的结果进行辨别
def filename_group(item_info, all_key, file_info):
    """
    :param items:
    :param all_key:
    :param file_info:
    :return:
    """
    if isinstance(item_info, list):
        # 判断value是否为列表
        for item in item_info:
            key, file_info = filename_group(item, all_key, file_info)
    elif isinstance(item_info, dict):
        for key, value in item_info.items():
            all_key += "/" + key
            key, file_info = filename_group(value, all_key, file_info)
    else:
        # 否则为字符串
        file_info[all_key] = item_info
    return all_key, file_info


def name_address(info, areas_content):
    """

    :param info: 传入正则匹配的规则
    :return:
    """
    name_list = []
    try:
        names = info.get("姓名")
        # address=info.get("地址")
        # if address:
        # for addr in address:
        # if len(addr) == 1:
        # pass
        # else:
        # address_list.append(addr)
        # info["地址"]=address_list
        if names:
            for name in names:
                if len(name) == 1:
                    pass
                else:
                    if name not in areas_content:
                        name_list.append(name)
            if name_list:
                info["姓名"] = name_list
            # if address:
            # info["地址"].extend(n_a)
            # else:
            # info["地址"]=[]
            # info["地址"].extend(n_a)
        return info
    except:
        return info


# 处理word modify rzc 2023/1/13
def word_con(filename, content):
    doc = docx.Document(filename)
    # 文件可能存在表格
    t_list = doc.tables

    for t in t_list:
        for i in range(len(t.rows)):
            for j in range(len(t.columns)):
                content += '"' + str(t.cell(i, j).text) + '"' + ","
    # 存在的段落
    gh_list = doc.paragraphs
    for paragraph in gh_list:
        content += paragraph.text
    return content


# 处理excel modify rzc 2023/1/13
def excel_con(zf, filename, content):
    global workbook
    zf.close()
    excel = False
    try:
        workbook = openpyxl.load_workbook(filename=filename)
        excel = True

    except:
        newname = filename + ".xlsx"
        if not os.path.exists(newname):
            os.rename(filename, newname)
        try:
            workbook = openpyxl.load_workbook(filename=newname)
            excel = True
        except:
            excel = False
        os.rename(newname, filename)
    # 获取所有表
    if excel:
        sheetnames = workbook.sheetnames
        # 对每个表都取数据
        for sheet in sheetnames:
            for i in workbook[sheet].rows:
                for j in i:
                    content += '"' + str(j.value) + '"' + ","
    return content


# add by rzc ppt 2023/1/13 处理excel
def ppt_con(filename, content):
    # 实例化ppt对象
    try:
        prs = Presentation(filename)
        results = []
        for slide in prs.slides:  # 幻灯片
            for shape in slide.shapes:  # 幻灯片内容
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        part = []
                        for run in paragraph.runs:
                            part.append(run.text)
                        results.append("".join(part))
        # 删除空得
        for r in results:
            if r == "":
                continue
            else:
                content += '"' + r + '"' + "\n"
        return content
    except:
        pass


# add by rzc name and address 2023/2/3 识别姓名地址
def identify_per_adr_name(text):
    """
    识别姓名
    :param text:
    :return:
    """
    try:
        words = pseg.cut(text, use_paddle=True)  # 词性标注，标注句子分词后每个词的词性
        name_jieba = []  # 识别出来得姓名
        address_jieba = []
        for pair_word in list(words):  # 遍历字符串
            # 分割字符串
            # print(list(pair_word))
            if list(pair_word)[1] == "nr" or list(pair_word)[1] == "PER":
                name_jieba.append(list(pair_word)[0])
            if list(pair_word)[1] == 'ns':
                address_jieba.append(list(pair_word)[0])
        return name_jieba, address_jieba
    except Exception as e:
        raise (e)


# add by rzc 对纳税人统一识别号进行规则检测2023/2/9
def check_social(code):
    # 统一代码由十八位的阿拉伯数字或大写英文字母（不使用I、O、Z、S、V）组成
    valid_char = '0123456789ABCDEFGHJKLMNPQRTUWXY'
    invlid_char = 'IOZSV'
    modulus = 31
    if code is None or code.isspace():
        return False
    # 位数校验
    if len(code) != 18:
        return False
    codeu = code.upper()
    # 含有不使用的大写英文字母
    for i in invlid_char:
        if i in codeu:
            return False
    # 不在有效的数字范围内
    for i in code:
        if not (i in valid_char):
            return False
    # 前两位登记管理部门代码和机构类别代码
    d_o_code = ['11', '12', '13', '19', '21', '29', '31', '32', '33', '34', '35', '39',
                '41', '49', '51', '52', '53', '59', '61', '62', '69', '71', '72', '79',
                '81', '89', '91', '92', '93', 'A1', 'A9', 'N1', 'N2', 'N3', 'N9', 'Y1']
    if not (code[0:2] in d_o_code):
        return False
    # 3-8位登记管理机关行政区划码先不检测了
    # 9-17位主体标识码（组织机构代码）先不拿出来校验了
    # 第18位校验码，这部分是关键
    # 第i位置对应的加权因子
    wi = [1, 3, 9, 27, 19, 26, 16, 17, 20, 29, 25, 13, 8, 24, 10, 30, 28]
    # 第i位置上的字符对应值
    corres_value = {'0': 0, '1': 1, '2': 2, '3': 3, '4': 4, '5': 5, '6': 6, '7': 7, '8': 8, '9': 9, 'A': 10,
                    'B': 11, 'C': 12, 'D': 13, 'E': 14, 'F': 15, 'G': 16, 'H': 17, 'J': 18, 'K': 19, 'L': 20,
                    'M': 21, 'N': 22, 'P': 23, 'Q': 24, 'R': 25, 'T': 26, 'U': 27, 'W': 28, 'X': 29, 'Y': 30}
    # 前17位字符位置序号i相对应的各个位置上的字符值
    ci = []
    for v in code[0:17]:
        j = corres_value.get(v)
        ci.append(j)
    # 计算与字符位置序号i相对应的乘积
    ciwi = []
    for i in range(len(ci)):
        ciwi.append(ci[i] * wi[i])
    # 计算级数之和
    ciwi_sum = sum(ciwi)
    # 级数之和求余,求出校验码字符值
    c18 = modulus - ciwi_sum % modulus
    # 查出校验码字符
    if c18 == modulus:
        c18 = 0
    check_code = None
    for k, v in corres_value.items():
        if v == c18:
            check_code = k
            break
    if check_code == code[17]:
        return True
    else:
        return False


# add by rzc 对银行卡号进行规则检测 2023/2/9
def luhn_valid(card_num):
    try:
        s = 0
        card_num_length = len(card_num)
        for _ in range(1, card_num_length + 1):
            t = int(card_num[card_num_length - _])
            if _ % 2 == 0:
                t *= 2
                s += t if t < 10 else t % 10 + t // 10
            else:
                s += t
        return s % 10 == 0
    except:
        return False


# add by rzc 对银行卡号进行校验 2023/4/7
def luhn(card_number):
    try:
        if not card_number.isdigit():
            return False
        if not 16 <= len(card_number) <= 19:
            return False
        digits = list(map(int, card_number))
        check_digit = digits.pop()
        digits.reverse()
        checksum = 0
        for i, digit in enumerate(digits):
            if i % 2 == 0:
                digit *= 2
                if digit > 9:
                    digit -= 9
            checksum += digit
        return (checksum + check_digit) % 10 == 0
    except:
        return False


# add by rzc 对身份证号进行规则校验 2023/2/15
def id_validators(card_num):
    an = validator.is_valid(card_num)
    return an


# 发票代码校验规则 2023/3/16
def check_invoice_code(code):
    """
    发票代码由12位数字组成，不能含有字母和特殊字符。

    第1位到第10位是发票的编码区，第11位和第12位是校验码。

    校验码的计算方法是将发票编码区的每一位数字乘以对应的权重，然后将乘积相加，得到一个结果。将结果除以11，得到一个余数。
    如果余数为0，则校验码为1；如果余数为1，则校验码为0；如果余数大于1，则校验码为11减去余数。
    :param code:
    :return:
    """
    if len(code) != 12:
        return False
    if not code.isdigit():  # 判断是否全为数字 如果不是返回False
        return False
    weights = [3, 7, 9, 10, 5, 8, 4, 2]
    total = sum(int(code[i]) * weights[i] for i in range(8))
    check_code = (11 - (total % 11)) % 11
    return check_code == int(code[10:])


# 发票代码校验规则 2023/3/18
def validate_invoice_code(code):
    # 校验码表
    check_codes = {0: 0, 1: 1, 2: 2, 3: 3, 4: 4, 5: 5, 6: 6, 7: 7, 8: 8, 9: 9, 10: 'X'}
    # 系数表
    factors = [3, 7, 9, 10, 5, 8, 4, 2, 1]
    # 计算校验和
    sum = 0
    for i in range(9):
        sum += int(code[i]) * factors[i]
    # 计算校验码
    check_code = check_codes[sum % 11]
    # 比较校验码和发票代码的第10位数字
    if str(check_code) == code[9]:
        return True
    else:
        return False


def check_invoice_code1(invoice_code):
    if not isinstance(invoice_code, str):
        return False
    if len(invoice_code) != 10:
        return False
    factors = [3, 7, 9, 10, 5, 8, 4, 2, 1]
    total = 0
    for i in range(9):
        if not invoice_code[i].isdigit():
            return False
        total += int(invoice_code[i]) * factors[i]
    remainder = total % 11
    if remainder == 0:
        check_code = '0'
    elif remainder == 1:
        check_code = 'X'
    else:
        check_code = str(11 - remainder)
    return check_code == invoice_code[9]


# 读取excel表格中存在行数及表单数 2023-6-30
def read_row(filepath, filename):
    """
    openpyxl.load_workbook()参数

    :param filepath:
    :param filename:
    :return:
    """
    global sheet_names, wb, workbook
    excel = False
    if filename.endswith(".xlsx"):
        # 读取excel表
        try:
            wb = openpyxl.load_workbook(filepath, read_only=True)
            excel = True
        except:
            newname = filepath + ".xlsx"
            if not os.path.exists(newname):
                os.rename(filepath, newname)
            try:
                wb = openpyxl.load_workbook(newname, read_only=True)
                excel = True
            except:
                excel = False
                pass
            # 转化为源文件
            os.rename(newname, filepath)
        if excel:
            # os.rename(filepath, newname)
            # 获取总共表单名称
            # 遍历每个工作表并获取行号
            sheet_names = wb.sheetnames
            row_counts = 0
            found = False
            # 循环表单信息
            for sheet_name in sheet_names:
                sheet = wb[sheet_name]
                num_rows = sheet.max_row
                if num_rows == 1:
                    # ws=wb.get_sheet_by_name(sheet_name)
                    for r in sheet.rows:
                        for c in r:
                            res = c.value
                            if res:
                                found = True
                                break
                    if found == True:
                        row_counts += num_rows
                else:
                    row_counts += num_rows
                # print(f"工作表'{sheet_name}'的行号:{num_rows}")
            # (row_counts)
            # 获取工作表的总数
            num_sheets = len(sheet_names)
            wb.close()
            # print(f'总共有{num_sheets}个工作表')
            return row_counts, num_sheets
    elif filename.endswith(".xls"):
        # 打开excel表
        try:
            workbook = xlrd.open_workbook(filepath, on_demand=False)
            excel = True
        except:
            newname = filepath + ".xls"
            if not os.path.exists(newname):
                os.rename(filepath, newname)
            try:
                workbook = xlrd.open_workbook(newname, on_demand=False)
                excel = True
            except:
                excel = False
            os.rename(newname, filepath)
        if excel:
            # 获取所有的sheetnames表
            sheets = workbook.sheet_names()
            row_counts = 0
            num_sheets = len(sheets)
            # 遍历工作表
            found = False
            for sheet_name in sheets:
                sheet = workbook.sheet_by_name(sheet_name)  # 选择当前sheets
                num_rows = sheet.nrows
                if num_rows == 1:
                    # 判断这一行是否为空行
                    for row_index in range(num_rows):
                        row = sheet.row(row_index)
                        for cell in row:
                            if cell.value:
                                found = True
                                break
                        if found:
                            break
                    if found:
                        row_counts += num_rows
                else:
                    row_counts += num_rows
            workbook.release_resources()
            return row_counts, num_sheets
    else:
        return 0, 0


if __name__ == '__main__':
    an = check_invoice_code1("4100194130")
    if an:
        print("这是银行卡号")
    else:
        print("不是")
